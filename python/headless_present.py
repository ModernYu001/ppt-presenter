#!/usr/bin/env python3
"""Headless PPT presenter: PPTX → slide images + per-slide narration audio → stitched video.

Usage:
    python3 headless_present.py input.pptx --out output.mp4 [--voice zh-CN-XiaoxiaoNeural]

Workflow:
    1. Parse PPTX, extract text + notes per slide
    2. Render each slide to PNG (python-pptx shapes → Pillow canvas)
    3. Generate narration text per slide via model API
    4. Generate TTS audio per slide via Edge TTS
    5. Stitch slides + audio into MP4 video with ffmpeg
"""

from __future__ import annotations

import argparse
import asyncio
import json
import os
import shutil
import subprocess
import sys
import tempfile
import textwrap
import time
import urllib.request
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import edge_tts

# --- Config ---
MODEL_BASE_URL = os.getenv("MODEL_BASE_URL", "http://54.179.184.88:8317/v1")
MODEL_API_KEY = os.getenv("MODEL_API_KEY", os.getenv("MODBTC_OPENAI_API_KEY", "modernyu002@key"))
MODEL_ID = os.getenv("MODEL_ID", "gpt-5.2")
MODEL_TIMEOUT = int(os.getenv("MODEL_TIMEOUT", "90"))
FFMPEG_BIN = os.getenv("FFMPEG_BIN", os.path.expanduser("~/bin/ffmpeg"))
SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080
DEFAULT_VOICE = "zh-CN-XiaoxiaoNeural"


# --- Step 1: Parse PPTX ---
def parse_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide:
                parts = []
                for shape in slide.notes_slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = shape.text.strip()
                        if t:
                            parts.append(t)
                notes = "\n".join(parts)
        except Exception:
            pass
        slides.append({
            "index": idx,
            "text": "\n".join([t for t in texts if t]),
            "notes": notes,
        })
    return slides


# --- Step 2: Render slide to image ---
def _try_load_font(size: int):
    """Try to load a CJK-capable font."""
    candidates = [
        "/usr/share/fonts/google-noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJKsc-Regular.otf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/liberation-sans/LiberationSans-Regular.ttf",
        "/usr/share/fonts/dejavu-sans-fonts/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
    return ImageFont.load_default()


def render_slide_image(slide_data: dict, out_path: str, width: int = SLIDE_WIDTH, height: int = SLIDE_HEIGHT):
    """Render slide text onto a clean image."""
    img = Image.new("RGB", (width, height), color=(24, 28, 44))
    draw = ImageDraw.Draw(img)

    title_font = _try_load_font(48)
    body_font = _try_load_font(32)
    idx_font = _try_load_font(20)

    # Slide number
    draw.text((width - 80, height - 40), f"{slide_data['index']}", fill=(100, 120, 160), font=idx_font)

    text = slide_data.get("text", "")
    lines = text.split("\n")
    title = lines[0] if lines else f"Slide {slide_data['index']}"
    body_lines = lines[1:] if len(lines) > 1 else []

    # Title
    y = 80
    # Word wrap title
    wrapped_title = textwrap.fill(title, width=40)
    for line in wrapped_title.split("\n"):
        draw.text((100, y), line, fill=(230, 235, 245), font=title_font)
        y += 60

    # Separator
    y += 20
    draw.line([(100, y), (width - 100, y)], fill=(0, 230, 138), width=2)
    y += 30

    # Body
    for bline in body_lines:
        if not bline.strip():
            y += 20
            continue
        wrapped = textwrap.fill(bline, width=55)
        for wl in wrapped.split("\n"):
            if y > height - 100:
                draw.text((100, y), "...", fill=(180, 190, 210), font=body_font)
                break
            draw.text((100, y), wl, fill=(200, 210, 225), font=body_font)
            y += 44
        if y > height - 100:
            break

    img.save(out_path, "PNG")
    return out_path


# --- Step 3: Generate narration via model ---
def generate_narration(slide_data: dict) -> str:
    prompt = (
        f"请为第 {slide_data['index']} 页 PPT 生成适合现场演讲的中文讲稿。\n"
        f"要求：口语化、流畅、不要机械复述；"
        "如果备注区有信息，优先吸收备注区意图；如果页内信息较少，可以做合理过渡。\n\n"
        f"【页面文字】\n{slide_data.get('text') or '(无)'}\n\n"
        f"【备注】\n{slide_data.get('notes') or '(无)'}\n\n"
        "只输出最终讲稿正文，不要标注页码或标题前缀。"
    )

    payload = json.dumps({
        "model": MODEL_ID,
        "messages": [
            {"role": "system", "content": "你是一位专业演讲者。生成自然、清晰的中文演讲稿，适合直接朗读。"},
            {"role": "user", "content": prompt},
        ],
        "temperature": 0.7,
    }).encode("utf-8")

    req = urllib.request.Request(
        f"{MODEL_BASE_URL}/chat/completions",
        data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {MODEL_API_KEY}",
        },
    )

    with urllib.request.urlopen(req, timeout=MODEL_TIMEOUT) as r:
        data = json.loads(r.read().decode("utf-8"))

    return data["choices"][0]["message"]["content"].strip()


# --- Step 4: TTS per slide ---
async def generate_tts(text: str, out_path: str, voice: str = DEFAULT_VOICE):
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(out_path)
    return out_path


# --- Step 5: Get audio duration ---
def get_audio_duration(path: str) -> float:
    """Get audio duration in seconds using ffprobe or ffmpeg."""
    ffprobe = FFMPEG_BIN.replace("ffmpeg", "ffprobe")
    if not os.path.exists(ffprobe):
        ffprobe = FFMPEG_BIN  # fallback
    try:
        result = subprocess.run(
            [FFMPEG_BIN, "-i", path, "-f", "null", "-"],
            capture_output=True, text=True, timeout=30,
        )
        # Parse duration from stderr
        for line in result.stderr.split("\n"):
            if "Duration:" in line:
                parts = line.split("Duration:")[1].split(",")[0].strip()
                h, m, s = parts.split(":")
                return float(h) * 3600 + float(m) * 60 + float(s)
    except Exception:
        pass
    return 10.0  # fallback


# --- Step 6: Stitch video ---
def stitch_video(slide_items: list[dict], out_path: str):
    """Stitch slide images + audio into a single MP4 video.

    Each slide shows for the duration of its audio.
    """
    tmpdir = tempfile.mkdtemp(prefix="ppt_video_")
    concat_file = os.path.join(tmpdir, "concat.txt")
    segment_files = []

    try:
        for item in slide_items:
            img_path = item["image"]
            audio_path = item["audio"]
            duration = get_audio_duration(audio_path)
            seg_path = os.path.join(tmpdir, f"seg_{item['index']:03d}.mp4")

            # Create video segment: still image + audio
            subprocess.run([
                FFMPEG_BIN,
                "-loop", "1",
                "-i", img_path,
                "-i", audio_path,
                "-c:v", "libx264",
                "-tune", "stillimage",
                "-c:a", "aac",
                "-b:a", "192k",
                "-pix_fmt", "yuv420p",
                "-shortest",
                "-t", str(duration + 0.5),  # small buffer
                "-y", seg_path,
            ], capture_output=True, timeout=120)

            if os.path.exists(seg_path):
                segment_files.append(seg_path)

        # Concat all segments
        with open(concat_file, "w") as f:
            for seg in segment_files:
                f.write(f"file '{seg}'\n")

        subprocess.run([
            FFMPEG_BIN,
            "-f", "concat",
            "-safe", "0",
            "-i", concat_file,
            "-c", "copy",
            "-y", out_path,
        ], capture_output=True, timeout=300)

    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)

    return out_path


# --- Main ---
def main():
    ap = argparse.ArgumentParser(description="Headless PPT presenter → video")
    ap.add_argument("pptx", help="Input PPTX file")
    ap.add_argument("--out", "-o", default=None, help="Output MP4 path")
    ap.add_argument("--voice", "-v", default=DEFAULT_VOICE, help="Edge TTS voice")
    ap.add_argument("--skip-narration", action="store_true", help="Use slide text directly instead of model narration")
    ap.add_argument("--workdir", default=None, help="Working directory for intermediates")
    args = ap.parse_args()

    pptx_path = os.path.abspath(args.pptx)
    if not os.path.exists(pptx_path):
        print(f"ERROR: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)

    out_path = args.out or os.path.splitext(pptx_path)[0] + ".mp4"
    out_path = os.path.abspath(out_path)

    workdir = args.workdir or tempfile.mkdtemp(prefix="ppt_present_")
    os.makedirs(workdir, exist_ok=True)

    print(f"[1/5] Parsing PPTX: {pptx_path}")
    slides = parse_pptx(pptx_path)
    print(f"  → {len(slides)} slides found")

    print(f"[2/5] Rendering slide images...")
    for slide in slides:
        img_path = os.path.join(workdir, f"slide_{slide['index']:03d}.png")
        render_slide_image(slide, img_path)
        slide["image"] = img_path
        print(f"  → slide {slide['index']} rendered")

    print(f"[3/5] Generating narration...")
    for slide in slides:
        if args.skip_narration:
            slide["narration"] = slide.get("text") or f"第{slide['index']}页"
        else:
            try:
                slide["narration"] = generate_narration(slide)
                print(f"  → slide {slide['index']} narration OK ({len(slide['narration'])} chars)")
            except Exception as e:
                print(f"  ⚠ slide {slide['index']} narration failed: {e}, using text fallback")
                slide["narration"] = slide.get("text") or f"第{slide['index']}页"

    print(f"[4/5] Generating TTS audio...")
    for slide in slides:
        audio_path = os.path.join(workdir, f"slide_{slide['index']:03d}.mp3")
        asyncio.run(generate_tts(slide["narration"], audio_path, voice=args.voice))
        slide["audio"] = audio_path
        print(f"  → slide {slide['index']} audio OK")

    print(f"[5/5] Stitching video → {out_path}")
    stitch_video(slides, out_path)

    if os.path.exists(out_path):
        size_mb = os.path.getsize(out_path) / (1024 * 1024)
        print(f"\n✅ Done! Output: {out_path} ({size_mb:.1f} MB)")
    else:
        print(f"\n❌ Failed to create {out_path}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
