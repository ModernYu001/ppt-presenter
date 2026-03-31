#!/usr/bin/env python3
"""Create a small test PPTX for headless_present.py testing."""
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "PPT 自动演示系统"
slide1.placeholders[1].text = "Edge TTS + 模型解说 + 视频合成"
notes1 = slide1.notes_slide
notes1.notes_text_frame.text = "这是系统测试的第一页，介绍项目概览。"

# Slide 2
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "核心功能"
slide2.placeholders[1].text = (
    "1. 自动解析 PPTX 内容\n"
    "2. AI 生成逐页演讲稿\n"
    "3. Edge TTS 语音合成\n"
    "4. 自动翻页视频输出"
)
notes2 = slide2.notes_slide
notes2.notes_text_frame.text = "重点强调端到端自动化，无需人工干预。"

# Slide 3
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "谢谢！"
slide3.placeholders[1].text = "欢迎提问与交流"

out = "/home/ec2-user/.openclaw/workspace/ppt-auto-presenter/test_deck.pptx"
prs.save(out)
print(f"Created: {out}")
