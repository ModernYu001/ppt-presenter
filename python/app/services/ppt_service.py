from pathlib import Path

from pptx import Presentation


class PptService:
    def load_presentation(self, path: str):
        ppt_path = Path(path)
        if not ppt_path.exists():
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")
        prs = Presentation(str(ppt_path))
        slides = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
            notes_text = ""
            try:
                if slide.has_notes_slide and slide.notes_slide:
                    parts = []
                    for shape in slide.notes_slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text = shape.text.strip()
                            if text:
                                parts.append(text)
                    notes_text = "\n".join(parts)
            except Exception:
                notes_text = ""
            slides.append(
                {
                    "index": idx,
                    "text": "\n".join([t for t in texts if t]),
                    "notes": notes_text,
                }
            )
        return slides
